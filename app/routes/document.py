from fastapi import APIRouter, Depends, HTTPException, UploadFile, File
from fastapi.responses import JSONResponse
from sqlalchemy.orm import Session
from typing import List
import io

import PyPDF2
from pptx import Presentation
from docx import Document as DocxDocument

from app.database import get_db
from app.models.document import Document
from app.models.project import Project
from app.schemas.document import DocumentCreate, DocumentResponse, DocumentUpdate

router = APIRouter(
    prefix="/documents",
    tags=["Documents"]
)

# ======================================
# EXTRACT TEXT FROM FILE
# ======================================
@router.post("/extract-text/")
async def extract_text(file: UploadFile = File(...)):

    try:
        print("1 - Extract route hit")

        file_bytes = await file.read()

        if not file_bytes:
            raise HTTPException(status_code=400, detail="Empty file")

        if len(file_bytes) > 10 * 1024 * 1024:
            raise HTTPException(status_code=400, detail="File too large (Max 10MB)")

        filename = file.filename.lower()
        content = ""
        file_stream = io.BytesIO(file_bytes)

        print("2 - File loaded into memory")

        # ================= TXT =================
        if filename.endswith(".txt"):
            content = file_bytes.decode("utf-8", errors="ignore")

        # ================= PDF =================
        elif filename.endswith(".pdf"):
            pdf_reader = PyPDF2.PdfReader(file_stream)

            if pdf_reader.is_encrypted:
                raise HTTPException(status_code=400, detail="Encrypted PDF not supported")

            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    content += text

        # ================= PPTX =================
        elif filename.endswith(".pptx"):
            presentation = Presentation(file_stream)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += shape.text + "\n"

        # ================= DOCX =================
        elif filename.endswith(".docx"):
            doc = DocxDocument(file_stream)
            for para in doc.paragraphs:
                content += para.text + "\n"

        # ================= DOC (Legacy) =================
        elif filename.endswith(".doc"):
            raise HTTPException(
                status_code=400,
                detail="Legacy .doc files not supported. Please convert to .docx."
            )

        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")

        print("3 - Processing complete")

        if not content.strip():
            raise HTTPException(status_code=400, detail="No readable text found in file.")

        return JSONResponse(content={"content": content})

    except HTTPException:
        raise

    except Exception as e:
        print("ERROR:", str(e))
        raise HTTPException(status_code=500, detail=str(e))


# ======================================
# GET DOCUMENTS
# ======================================
@router.get("/", response_model=List[DocumentResponse])
def get_documents(project_id: int = None, db: Session = Depends(get_db)):
    query = db.query(Document)

    if project_id is not None:
        query = query.filter(Document.project_id == project_id)

    return query.all()


# ======================================
# CREATE DOCUMENT
# ======================================
@router.post("/", response_model=DocumentResponse)
def create_document(document: DocumentCreate, db: Session = Depends(get_db)):

    project = db.query(Project).filter(Project.id == document.project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    db_document = Document(
        title=document.title,
        content=document.content,
        project_id=document.project_id
    )

    db.add(db_document)
    db.commit()
    db.refresh(db_document)

    return db_document


# ======================================
# UPDATE DOCUMENT
# ======================================
@router.put("/{document_id}", response_model=DocumentResponse)
def update_document(document_id: int, data: DocumentUpdate, db: Session = Depends(get_db)):

    document = db.query(Document).filter(Document.id == document_id).first()

    if not document:
        raise HTTPException(status_code=404, detail="Document not found")

    for key, value in data.dict(exclude_unset=True).items():
        setattr(document, key, value)

    db.commit()
    db.refresh(document)

    return document


# ======================================
# DELETE DOCUMENT
# ======================================
@router.delete("/{document_id}")
def delete_document(document_id: int, db: Session = Depends(get_db)):

    document = db.query(Document).filter(Document.id == document_id).first()

    if not document:
        raise HTTPException(status_code=404, detail="Document not found")

    db.delete(document)
    db.commit()

    return {"message": "Document deleted successfully"}